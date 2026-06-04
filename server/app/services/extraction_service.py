# hf_final_deploy/app/services/extraction_service.py
import re, io, html, asyncio, base64
from typing import Optional, List, Dict


def _extract_youtube_id(url: str):
    patterns = [
        r"(?:youtube\.com/watch\?(?:.*&)?v=|youtu\.be/|youtube\.com/embed/|youtube\.com/shorts/|youtube\.com/v/)([0-9A-Za-z_-]{11})",
        r"(?:v=|vi=|v%3D)([0-9A-Za-z_-]{11})",
        r"^([0-9A-Za-z_-]{11})$",
    ]
    for pattern in patterns:
        m = re.search(pattern, url.strip())
        if m:
            return m.group(1)
    return None


def _items_to_text(data) -> str:
    parts = []
    for item in data:
        if isinstance(item, dict):
            parts.append(item.get("text", ""))
        else:
            text = getattr(item, "text", None)
            parts.append(text if text is not None else str(item))
    return " ".join(p.strip() for p in parts if p.strip())


async def _transcript_api(video_id: str) -> str:
    from youtube_transcript_api import YouTubeTranscriptApi, TranscriptsDisabled, NoTranscriptFound
    loop = asyncio.get_event_loop()

    def _fetch():
        try:
            return _items_to_text(YouTubeTranscriptApi.get_transcript(video_id, languages=["en", "en-US", "en-GB", "en-CA"]))
        except Exception:
            pass
        try:
            tlist = YouTubeTranscriptApi.list_transcripts(video_id)
            for finder in [
                lambda: tlist.find_manually_created_transcript(["en", "en-US", "en-GB"]),
                lambda: tlist.find_generated_transcript(["en", "en-US", "en-GB"]),
            ]:
                try:
                    return _items_to_text(finder().fetch())
                except Exception:
                    pass
            try:
                t = next(iter(tlist))
                if t.is_translatable:
                    return _items_to_text(t.translate("en").fetch())
                return _items_to_text(t.fetch())
            except Exception:
                pass
        except (TranscriptsDisabled, NoTranscriptFound) as e:
            raise ValueError(f"Transcripts disabled: {e}")
        except Exception as e:
            raise ValueError(f"Transcript API error: {e}")
        raise ValueError("Could not retrieve transcript via API.")

    return await loop.run_in_executor(None, _fetch)


async def _transcript_proxy(video_id: str) -> str:
    import requests as req
    r = req.get(f"https://yt-transcript.vercel.app/api/transcript?v={video_id}", timeout=12)
    if r.status_code == 200:
        data = r.json()
        if isinstance(data, list) and data:
            return _items_to_text(data)
    raise ValueError("Vercel proxy returned no data.")


async def _transcript_scrape(video_id: str) -> str:
    import requests as req
    r = req.get(
        f"https://youtubetranscript.com/?server_vid2={video_id}",
        timeout=12,
        headers={"User-Agent": "Mozilla/5.0"},
    )
    if r.status_code == 200:
        matches = re.findall(r"<text[^>]*>(.*?)</text>", r.text, re.DOTALL)
        if matches:
            return html.unescape(" ".join(m.strip() for m in matches))
    raise ValueError("Scraper returned no data.")


def _extract_pdf(content: bytes) -> str:
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page in pdf.pages:
                t = page.extract_text(x_tolerance=3, y_tolerance=3)
                if t:
                    parts.append(t)
        text = "\n".join(parts)
        if len(text.strip()) > 50:
            print(f"[LUMINA] PDF via pdfplumber: {len(text)} chars")
            return text
    except Exception as e:
        print(f"[LUMINA] pdfplumber failed: {e}")
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(content))
        parts = [p.extract_text() for p in reader.pages if p.extract_text()]
        text = "\n".join(parts)
        if len(text.strip()) > 10:
            print(f"[LUMINA] PDF via PyPDF2: {len(text)} chars")
            return text
    except Exception as e:
        print(f"[LUMINA] PyPDF2 failed: {e}")
    raise ValueError("Could not extract text from PDF. File may be scanned or password-protected.")


def _extract_docx(content: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(content))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    result = "\n".join(parts)
    print(f"[LUMINA] DOCX extracted: {len(result)} chars")
    return result


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if row_text:
                        slide_parts.append(row_text)
        if len(slide_parts) > 1:
            parts.extend(slide_parts)
    result = "\n".join(parts)
    print(f"[LUMINA] PPTX extracted: {len(result)} chars, {len(prs.slides)} slides")
    return result


def _extract_txt(content: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            text = content.decode(enc)
            print(f"[LUMINA] TXT decoded ({enc}): {len(text)} chars")
            return text
        except (UnicodeDecodeError, LookupError):
            continue
    raise ValueError("Could not decode text file. Please use UTF-8 encoding.")


def _extract_pdf_figures(content: bytes, max_images: int = 8) -> List[Dict]:
    figures: List[Dict] = []
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page_num, page in enumerate(pdf.pages[:15]):
                if len(figures) >= max_images or not page.images:
                    continue
                try:
                    # Sort images on the page by top-to-bottom and left-to-right
                    sorted_imgs = sorted(page.images, key=lambda x: (x.get("top", 0), x.get("x0", 0)))
                    for img in sorted_imgs[:5]:
                        if len(figures) >= max_images:
                            break
                        x0, top, x1, bottom = img.get("x0"), img.get("top"), img.get("x1"), img.get("bottom")
                        if x0 is None or top is None or x1 is None or bottom is None:
                            continue
                        
                        # Ensure bounding box coordinates are ordered and valid
                        x0_c = min(x0, x1)
                        x1_c = max(x0, x1)
                        top_c = min(top, bottom)
                        bottom_c = max(top, bottom)
                        
                        # Add a tiny padding to prevent borders from being clipped
                        padding = 2
                        x0_c = max(0, x0_c - padding)
                        top_c = max(0, top_c - padding)
                        x1_c = min(page.width, x1_c + padding)
                        bottom_c = min(page.height, bottom_c + padding)
                        
                        # Check for degenerate bounding box
                        if (x1_c - x0_c) < 30 or (bottom_c - top_c) < 30:
                            continue
                            
                        # If the image covers more than 95% of the page width and height, it's likely a full-page scan, not a figure
                        if (x1_c - x0_c) > 0.95 * page.width and (bottom_c - top_c) > 0.95 * page.height:
                            continue
                            
                        # Crop the page first - pdfplumber handles page rotation, CropBox translation, etc. natively
                        cropped_page = page.crop((x0_c, top_c, x1_c, bottom_c), strict=False)
                        cropped_image = cropped_page.to_image(resolution=150)
                        crop = cropped_image.original
                        
                        if crop.width < 50 or crop.height < 50:
                            continue
                            
                        buf = io.BytesIO()
                        crop.save(buf, format="PNG")
                        b64 = base64.b64encode(buf.getvalue()).decode("ascii")
                        figures.append({
                            "url": f"data:image/png;base64,{b64}",
                            "alt": f"Figure from page {page_num + 1}",
                            "caption": f"From your uploaded PDF (page {page_num + 1})",
                            "source": "upload",
                        })
                except Exception as e:
                    print(f"[LUMINA] PDF figure page {page_num + 1}: {e}")
    except Exception as e:
        print(f"[LUMINA] PDF figures failed: {e}")
    return figures


def _extract_pptx_figures(content: bytes, max_images: int = 8) -> List[Dict]:
    figures: List[Dict] = []
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(io.BytesIO(content))
        for slide_num, slide in enumerate(prs.slides, 1):
            if len(figures) >= max_images:
                break
            for shape in slide.shapes:
                if len(figures) >= max_images:
                    break
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        blob = shape.image.blob
                        ext = (shape.image.ext or "png").lower()
                        mime = "image/jpeg" if ext in ("jpg", "jpeg") else f"image/{ext}"
                        b64 = base64.b64encode(blob).decode("ascii")
                        figures.append({
                            "url": f"data:{mime};base64,{b64}",
                            "alt": f"Slide {slide_num} figure",
                            "caption": f"From your uploaded slides (slide {slide_num})",
                            "source": "upload",
                        })
                    except Exception as e:
                        print(f"[LUMINA] PPTX image slide {slide_num}: {e}")
    except Exception as e:
        print(f"[LUMINA] PPTX figures failed: {e}")
    return figures


class ExtractionService:

    @staticmethod
    def extract_figures_from_file(file_content: bytes, filename: str) -> List[Dict]:
        if not filename:
            return []
        name = filename.lower()
        if name.endswith(".pdf"):
            return _extract_pdf_figures(file_content)
        if name.endswith((".pptx", ".ppt")):
            return _extract_pptx_figures(file_content)
        return []

    @staticmethod
    def extract_youtube_id(url: str):
        return _extract_youtube_id(url)

    @staticmethod
    async def get_youtube_transcript(url: str) -> str:
        video_id = _extract_youtube_id(url)
        if not video_id:
            raise ValueError(
                "Invalid YouTube URL. Supported: youtube.com/watch?v=ID, youtu.be/ID, youtube.com/shorts/ID"
            )
        print(f"[LUMINA] Extracting transcript: {video_id}")

        # Tier 1: Official API
        try:
            text = await _transcript_api(video_id)
            if text and len(text.strip()) > 50:
                print(f"[LUMINA] Transcript via API: {len(text)} chars")
                return text
        except ValueError as e:
            err = str(e).lower()
            if "disabled" in err or "not found" in err:
                raise ValueError(
                    "Captions are disabled for this video. "
                    "Try pasting the transcript manually using the link below."
                )
            print(f"[LUMINA] API tier failed: {e}")
        except Exception as e:
            print(f"[LUMINA] API tier error: {e}")

        # Tier 2: Vercel proxy
        try:
            text = await _transcript_proxy(video_id)
            if text and len(text.strip()) > 50:
                print(f"[LUMINA] Transcript via proxy: {len(text)} chars")
                return text
        except Exception as e:
            print(f"[LUMINA] Proxy tier failed: {e}")

        # Tier 3: Scraper
        try:
            text = await _transcript_scrape(video_id)
            if text and len(text.strip()) > 50:
                print(f"[LUMINA] Transcript via scraper: {len(text)} chars")
                return text
        except Exception as e:
            print(f"[LUMINA] Scraper tier failed: {e}")

        raise ValueError(
            "Could not extract transcript. YouTube may have blocked the request. "
            "Please paste the transcript manually using the link below."
        )

    @staticmethod
    async def extract_text_from_file(file_content: bytes, filename: str) -> str:
        if not filename:
            raise ValueError("Filename required to determine file type.")
        name = filename.lower()
        print(f"[LUMINA] Extracting file: {filename} ({len(file_content)} bytes)")

        if name.endswith(".pdf"):
            return _extract_pdf(file_content)
        elif name.endswith((".docx", ".doc")):
            try:
                return _extract_docx(file_content)
            except Exception as e:
                raise ValueError(f"Could not read Word document: {e}")
        elif name.endswith((".pptx", ".ppt")):
            try:
                return _extract_pptx(file_content)
            except Exception as e:
                raise ValueError(f"Could not read PowerPoint file: {e}")
        elif name.endswith((".txt", ".md", ".markdown", ".rst", ".csv")):
            return _extract_txt(file_content)
        else:
            ext = filename.rsplit(".", 1)[-1].upper() if "." in filename else "unknown"
            raise ValueError(
                f"Unsupported file type: .{ext}. Supported: PDF, DOCX, PPTX, TXT, MD"
            )
